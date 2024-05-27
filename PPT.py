from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
import matplotlib.pyplot as plt
import pandas as pd

# Создание объекта презентации
prs = Presentation()

# Первая страница (Титульный лист)
slide_title = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_title.shapes.title
subtitle = slide_title.placeholders[1]
title.text = "Анализ данных пингвинов"
subtitle.text = "Имя: XXX\nДата: XX.XX.XXXX"

# Вторая страница (Описание набора данных)
slide_intro = prs.slides.add_slide(prs.slide_layouts[1])
title_intro = slide_intro.shapes.title
content_intro = slide_intro.placeholders[1]
title_intro.text = "Описание набора данных"
content_intro.text = (
    "Название набора данных: Пингвины\n"
    "Источник данных: библиотека Seaborn\n"
    "Описание данных:\n"
    "- Набор данных о пингвинах содержит три вида пингвинов: Adelie, Chinstrap и Gentoo.\n"
    "- Каждый пингвин характеризуется такими признаками, как длина клюва, глубина клюва, длина плавника, масса тела и пол.\n"
    "- Данные собраны на различных островах: Biscoe, Dream и Torgersen.\n\n"
    "Характеристики данных:\n"
    "- species: Вид пингвина\n"
    "- island: Остров\n"
    "- bill_length_mm: Длина клюва (мм)\n"
    "- bill_depth_mm: Глубина клюва (мм)\n"
    "- flipper_length_mm: Длина плавника (мм)\n"
    "- body_mass_g: Масса тела (г)\n"
    "- sex: Пол"
)

# Третья страница (Предобработка данных)
slide_preprocessing = prs.slides.add_slide(prs.slide_layouts[1])
title_preprocessing = slide_preprocessing.shapes.title
content_preprocessing = slide_preprocessing.placeholders[1]
title_preprocessing.text = "Предобработка данных"
content_preprocessing.text = (
    "Шаги предобработки данных:\n"
    "1. Проверка на пропущенные данные:\n"
    "   - Использование метода isnull().sum() из библиотеки pandas для проверки пропущенных значений в каждом столбце.\n"
    "   - Результаты показывают, что в некоторых столбцах есть пропущенные значения.\n"
    "2. Удаление пропущенных данных:\n"
    "   - Использование метода dropna() для удаления строк с пропущенными значениями.\n"
    "   - После обработки набор данных больше не содержит пропущенных значений.\n"
    "3. Проверка результатов обработки пропущенных данных:\n"
    "   - Повторная проверка пропущенных значений в каждом столбце для подтверждения их удаления.\n\n"
    "Код предобработки данных:\n"
    "import seaborn as sns\n"
    "import pandas as pd\n\n"
    "# Загрузка набора данных о пингвинах\n"
    "penguins = sns.load_dataset(\"penguins\")\n\n"
    "# Проверка на пропущенные данные\n"
    "print(\"Пропущенные данные:\")\n"
    "print(penguins.isnull().sum())\n\n"
    "# Удаление пропущенных данных\n"
    "penguins_cleaned = penguins.dropna()\n\n"
    "# Проверка результатов обработки пропущенных данных\n"
    "print(\"После обработки пропущенных данных:\")\n"
    "print(penguins_cleaned.isnull().sum())"
)

# Функция для добавления слайда с графиком и выводом
def add_slide_with_chart(prs, title_text, chart_func, fig_path, conclusion_text):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = title_text
    chart_func()
    plt.savefig(fig_path)
    plt.close()
    slide.shapes.add_picture(fig_path, Inches(1), Inches(1.5), width=Inches(8))
    textbox = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1.5))
    textbox.text = conclusion_text

penguins_cleaned = sns.load_dataset("penguins").dropna()

# Количество видов пингвинов
add_slide_with_chart(
    prs,
    "Количество видов пингвинов",
    lambda: sns.countplot(data=penguins_cleaned, x='species'),
    "penguin_species_count.png",
    "Вывод: Диаграмма количества видов пингвинов показывает, что пингвинов вида Adelie больше всего, а пингвинов вида Chinstrap меньше всего."
)

# Количество видов пингвинов на разных островах
add_slide_with_chart(
    prs,
    "Количество видов пингвинов на разных островах",
    lambda: sns.countplot(data=penguins_cleaned, x='island', hue='species'),
    "penguin_island_species_count.png",
    "Вывод: Диаграмма количества видов пингвинов на разных островах показывает, что пингвины вида Adelie распространены на всех островах, тогда как пингвины видов Chinstrap и Gentoo сконцентрированы на определенных островах."
)

# Распределение длины плавника
add_slide_with_chart(
    prs,
    "Распределение длины плавника",
    lambda: sns.histplot(data=penguins_cleaned, x='flipper_length_mm', hue='species', kde=True),
    "penguin_flipper_length_dist.png",
    "Вывод: Диаграмма распределения длины плавника показывает, что длина плавников у пингвинов вида Gentoo заметно больше, чем у других видов."
)

# Взаимосвязь длины и глубины клюва
add_slide_with_chart(
    prs,
    "Взаимосвязь длины и глубины клюва",
    lambda: sns.scatterplot(data=penguins_cleaned, x='bill_length_mm', y='bill_depth_mm', hue='species'),
    "penguin_bill_length_depth.png",
    "Вывод: Диаграмма взаимосвязи длины и глубины клюва показывает, что у разных видов пингвинов размеры клюва различаются."
)

# Линейная регрессия массы тела и длины плавника
add_slide_with_chart(
    prs,
    "Линейная регрессия массы тела и длины плавника",
    lambda: sns.lmplot(data=penguins_cleaned, x='body_mass_g', y='flipper_length_mm', hue='species'),
    "penguin_body_mass_flipper_length.png",
    "Вывод: Диаграмма линейной регрессии массы тела и длины плавника показывает, что между массой тела и длиной плавника у пингвинов существует положительная корреляция."
)

# Ящик с усами
add_slide_with_chart(
    prs,
    "Ящик с усами",
    lambda: sns.boxplot(data=penguins_cleaned, x='species', y='flipper_length_mm'),
    "penguin_flipper_length_box.png",
    "Вывод: Диаграмма распределения длины плавника показывает, что длина плавников у разных видов пингвинов существенно различается, и у пингвинов вида Gentoo она больше."
)

# Объединенная диаграмма
add_slide_with_chart(
    prs,
    "Объединенная диаграмма",
    lambda: sns.jointplot(data=penguins_cleaned, x='bill_length_mm', y='bill_depth_mm', hue='species', kind='scatter'),
    "penguin_jointplot.png",
    "Вывод: Объединенная диаграмма показывает, что взаимосвязь длины и глубины клюва у разных видов пингвинов заметно различается."
)

# Тепловая карта корреляции
add_slide_with_chart(
    prs,
    "Тепловая карта корреляции",
    lambda: sns.heatmap(penguins_cleaned.select_dtypes(include='number').corr(), annot=True, cmap='coolwarm', center=0),
    "penguin_correlation_heatmap.png",
    "Вывод: Тепловая карта корреляции показывает, что между массой тела и длиной плавника у пингвинов существует сильная положительная корреляция."
)

# Распределение массы тела у разных видов пингвинов
add_slide_with_chart(
    prs,
    "Распределение массы тела у разных видов пингвинов",
    lambda: sns.boxplot(data=penguins_cleaned, x='species', y='body_mass_g'),
    "penguin_body_mass_box.png",
    "Вывод: Диаграмма распределения массы тела показывает, что масса тела у пингвинов вида Gentoo существенно больше, чем у других видов."
)

# Общий вывод
slide_summary = prs.slides.add_slide(prs.slide_layouts[1])
title_summary = slide_summary.shapes.title
content_summary = slide_summary.placeholders[1]
title_summary.text = "Общий вывод"
content_summary.text = (
    "Анализ набора данных о пингвинах показал, что пингвинов вида Adelie больше всего, и они наиболее широко распространены.\n"
    "Пингвины вида Gentoo имеют заметно большую длину плавников и массу тела по сравнению с другими видами.\n"
    "У разных видов пингвинов длина и глубина клюва существенно различаются.\n"
    "Между массой тела и длиной плавника у пингвинов существует положительная корреляция.\n"
    "В целом, различные виды пингвинов имеют заметные различия в различных характеристиках, которые можно наглядно показать с помощью различных методов визуализации."
)

# Сохранение презентации
prs.save("penguin_analysis_presentation.pptx")
